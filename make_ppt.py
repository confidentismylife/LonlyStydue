"""生成组会演示 PPT — 火灾疏散行人轨迹预测 (白色底 + 详细版)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 配色方案（白色底） ──
BG   = RGBColor(0xff, 0xff, 0xff)  # 白底
TITLE= RGBColor(0x1a, 0x1a, 0x2e)  # 深色标题
TEXT = RGBColor(0x33, 0x33, 0x44)  # 正文
META = RGBColor(0x88, 0x88, 0x99)  # 次要文字
RED  = RGBColor(0xdc, 0x26, 0x26)  # 强调红
BLUE = RGBColor(0x25, 0x6d, 0xeb)  # 主色蓝
LBLUE= RGBColor(0xe8, 0xf0, 0xfe)  # 浅蓝底
GREEN= RGBColor(0x05, 0x96, 0x69)  # 完成绿
ORANGE=RGBColor(0xd9, 0x77, 0x06)  # 警告橙
LGRAY= RGBColor(0xf5, 0xf5, 0xf8)  # 卡片底
BORDER=RGBColor(0xe0, 0xe0, 0xe8)  # 边框

def add_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, size=16,
             color=TEXT, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tf

def add_shape(slide, left, top, w, h, fill_color=LGRAY, border_color=BORDER):
    shape = slide.shapes.add_shape(5, Inches(left), Inches(top),
                                   Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color; shape.line.width = Pt(0.5)
    return shape

def add_rounded_shape(slide, left, top, w, h, fill_color, text, size=11, color=TEXT, bold=False):
    """圆角矩形标签"""
    shape = slide.shapes.add_shape(
        5, Inches(left), Inches(top), Inches(w), Inches(h))  # 5 = rounded rect
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = 'Microsoft YaHei'
    return shape

def add_arrow(slide, left, top, w, h):
    """右箭头"""
    shape = slide.shapes.add_shape(
        33, Inches(left), Inches(top), Inches(w), Inches(h))  # 33 = right arrow
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape

def add_page_number(slide, num):
    add_text(slide, 12.2, 7.0, 0.8, 0.3, str(num), size=10, color=META, align=PP_ALIGN.RIGHT)

def add_section_line(slide, top):
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(top), Inches(2.0), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()

def add_card(slide, left, top, w, h, title, lines, icon="●"):
    """白色卡片+蓝色左边框"""
    shape = add_shape(slide, left, top, w, h, LGRAY, BORDER)
    add_text(slide, left+0.2, top+0.1, w-0.4, 0.35,
             f'{icon} {title}', size=14, color=BLUE, bold=True)
    body = '\n'.join(f'  {l}' for l in lines)
    add_text(slide, left+0.2, top+0.45, w-0.4, h-0.55, body, size=11, color=TEXT)

def add_data_table(slide, left, top, headers, rows, col_widths, highlight_row=None):
    """绘制数据表格"""
    n_cols = len(headers)
    # 表头
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        shape = add_shape(slide, left+sum(col_widths[:j]), top, w, 0.35, BLUE, BLUE)
        tf = shape.text_frame; p = tf.paragraphs[0]
        p.text = h; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10); p.font.color.rgb = BG; p.font.bold = True
    # 数据行
    for i, row in enumerate(rows):
        bg_color = LBLUE if highlight_row == i else (LGRAY if i%2==0 else BG)
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            shape = add_shape(slide, left+sum(col_widths[:j]), top+0.35+i*0.32, w, 0.32, bg_color)
            tf = shape.text_frame; p = tf.paragraphs[0]
            p.text = str(cell); p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10); p.font.color.rgb = RED if highlight_row==i else TEXT

# ══════════════════════════════════════════════════════════════
# Slide 1 — 封面
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
# 顶部装饰条
shape = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
# 左侧色块
shape = s.shapes.add_shape(1, Inches(0), Inches(2.8), Inches(5.0), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()
add_text(s, 1.0, 1.2, 11, 1.0, '火灾疏散场景下的行人轨迹预测', 40, TITLE, True)
add_text(s, 1.0, 2.4, 11, 0.6, '基于 Social-LSTM 的社交池化轨迹预测模型', 22, BLUE, False)
add_text(s, 1.0, 3.2, 11, 0.4, 'Social-LSTM: Human Trajectory Prediction in Crowded Spaces  |  CVPR 2016', 13, META)
add_text(s, 1.0, 5.5, 11, 0.4, '组 会 汇 报', 24, TITLE, True)
add_text(s, 1.0, 6.1, 11, 0.4, '2026 年 5 月', 14, META)
add_text(s, 1.0, 6.5, 11, 0.3, 'github.com/confidentismylife/LonlyStydue', 10, META)
add_page_number(s, 1)

# ══════════════════════════════════════════════════════════════
# Slide 2 — 目录
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 10, 0.7, '目  录', 32, TITLE, True)
add_section_line(s, 1.2)
items = [
    ('01', '研究背景与问题定义', '火灾疏散的痛点 & 轨迹预测的价值'),
    ('02', '数据集：ETH/UCY', '数据格式、场景统计、评估协议'),
    ('03', '数据预处理管道', '场景级加载、归一化、滑动窗口'),
    ('04', '模型架构：Social-LSTM', '社交池化、LSTMCell、编码-解码'),
    ('05', '训练策略与加速优化', '场景打包批处理、梯度管理、CPU 加速'),
    ('06', '实验结果与分析', 'ADE/FDE、模型对比、消融分析'),
    ('07', '代码工程架构', '模块化设计、文件职责、技术栈'),
    ('08', '未来工作与论文方向', '火灾特征、多模态、损失升级'),
]
for i, (num, title, desc) in enumerate(items):
    y = 1.5 + i * 0.68
    add_text(s, 1.5, y, 0.6, 0.4, num, 22, BLUE, True)
    add_text(s, 2.2, y, 5.0, 0.35, title, 16, TITLE, True)
    add_text(s, 2.2, y+0.32, 6.0, 0.3, desc, 11, META)
add_page_number(s, 2)

# ══════════════════════════════════════════════════════════════
# Slide 3 — 背景与问题
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '01  研究背景与问题定义', 28, TITLE, True)
add_section_line(s, 1.0)

add_card(s, 0.8, 1.3, 3.8, 2.4, '火灾疏散的现实痛点',
    ['火灾发生时人群恐慌，行为高度不确定',
     '传统模拟工具 (Pathfinder/FDS) 假设行人理性匀速行走',
     '烟雾/高温/能见度变化严重影响决策',
     '缺少数据驱动的方法来建模真实行人交互',
     '现有疏散方案忽视了个体间微观行为差异',
     '实时轨迹预测 → 动态引导 → 减少伤亡'], RED)

add_card(s, 5.0, 1.3, 3.8, 2.4, '轨迹预测的技术价值',
    ['输入：行人过去 3.2s (8帧) 的位置序列',
     '输出：未来 4.8s (12帧) 的运动轨迹',
     '识别拥堵瓶颈 → 提前预警',
     '为消防机器人/无人机提供避障路径',
     '智慧楼宇数字孪生的核心感知模块',
     '辅助应急指挥决策'], BLUE)

add_card(s, 9.2, 1.3, 3.5, 2.4, '核心研究问题',
    ['Q1: 如何建模行人间的社交交互？',
     '  → 社交池化 (Social Pooling)',
     'Q2: 如何让模型"看见"周围人？',
     '  → 空间网格 + 隐藏状态求和',
     'Q3: 如何高效训练（无 GPU）？',
     '  → 场景打包批处理 (3.4x 加速)'], ORANGE)

# 问题形式化
add_text(s, 0.8, 4.1, 11.5, 0.3, '▎问题形式化', 15, TITLE, True)
add_shape(s, 0.8, 4.5, 11.7, 1.6, LBLUE, BLUE)
add_text(s, 1.2, 4.6, 11, 1.4,
    '给定:  T_obs = [p₁ᵗ⁻⁷, p₁ᵗ⁻⁶, ..., p₁ᵗ⁻¹, p₁ᵗ]  × N 个行人    其中 p = (x, y)\n'
    '       ───────── 8 帧观测轨迹 ─────────\n\n'
    '目标:  T_pred = [p̂₁ᵗ⁺¹, p̂₁ᵗ⁺², ..., p̂₁ᵗ⁺¹²]                       预测未来 12 帧位置\n'
    '       ───────── 12 帧预测轨迹 ─────────\n\n'
    '约束:  考虑所有邻居 j 对行人 i 的影响 (j ≠ i)，即社交交互项 Sᵢᵗ = Pooling(hⱼᵗ⁻¹, pⱼᵗ − pᵢᵗ)',
    size=12, color=TEXT)
add_page_number(s, 3)

# ══════════════════════════════════════════════════════════════
# Slide 4 — 数据集
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '02  数据集：ETH / UCY', 28, TITLE, True)
add_section_line(s, 1.0)

add_text(s, 0.8, 1.2, 10, 0.3, '▎5 个场景统计', 14, TITLE, True)
add_data_table(s, 0.8, 1.6,
    ['场景', '文件', '行人数', '帧数', '帧范围', '密度(人/帧)', '场景类型'],
    [['ETH',   'biwi_eth.txt',     '360', '876',  '780-12,380',  '6.3',  '街道'],
     ['Hotel', 'biwi_hotel.txt',   '389', '1,168','0-18,060',    '5.6',  '酒店大堂'],
     ['Univ',  'students001/003',  '415', '444',  '0-4,430',     '49.1', '校园广场 ★'],
     ['Zara1', 'crowds_zara01.txt','148', '872',  '0-9,010',     '5.9',  '商业街'],
     ['Zara2', 'crowds_zara02.txt','204', '1,052','10-10,520',   '9.2',  '商业街']],
    [0.9, 2.1, 0.8, 0.7, 1.6, 1.2, 1.0],
    highlight_row=2)

add_card(s, 0.8, 3.7, 5.8, 1.8, '数据格式与预处理',
    ['文件格式: frame_id  ped_id  x  y   (Tab分隔, 无表头)',
     '示例:  780.0  1.0  8.46  3.59',
     '帧间隔: Δframe = 10 (25fps原视频 → 2.5fps采样)',
     '0.4 秒/帧 × 8 帧 = 3.2 秒观测    |    0.4 秒/帧 × 12 帧 = 4.8 秒预测',
     '归一化策略: 每条轨迹以最后一帧观测位置为原点 (标准做法)',
     'origin = obs[:,-1,:]; obs_norm = obs - origin; pred_norm = pred - origin'],
    BLUE)

add_card(s, 7.0, 3.7, 5.8, 1.8, '评估协议 (Leave-One-Out)',
    ['标准做法: 4 个场景训练 → 第 5 个场景测试 → 轮换 5 次取平均',
     'ADE (Average Displacement Error): 12 帧预测与真值的平均L2距离',
     '  ADE = (1/N) Σᵢ (1/12) Σₜ ||p̂ᵢᵗ − pᵢᵗ||₂',
     'FDE (Final Displacement Error): 仅最后一帧的L2距离',
     '  FDE = (1/N) Σᵢ ||p̂ᵢᵗ¹² − pᵢᵗ¹²||₂',
     '训练/验证/测试严格按时间划分: 前80%训练 + 后20%验证 + 独立测试集'],
    GREEN)

add_text(s, 0.8, 6.0, 11.5, 0.4,
    '▸ 数据来源: Social-STGCNN 仓库 (github.com/abduallahmohamed/Social-STGCNN) — 基于 ETH 和 UCY 原始数据重新划分',
    size=11, color=META)
add_page_number(s, 4)

# ══════════════════════════════════════════════════════════════
# Slide 5 — 数据预处理管道
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '03  数据预处理管道', 28, TITLE, True)
add_section_line(s, 1.0)

# 流程图
steps = [
    ('原始文件', '.txt 文件\nframe ped x y', LGRAY),
    ('load_scene_data()\n按帧分组', 'frames[frame_id]\n= {ped: (x,y)}', LBLUE),
    ('extract_scene_windows()\n滑动窗口', 'obs (P,8,2)\npred (P,12,2)', LBLUE),
    ('normalize_scene_window()\n轨迹归一化', '最后观测点\n→ 原点', LBLUE),
    ('打包批处理\npack_scenes()', '6个场景→1批\n(total_P, 8, 2)', BLUE),
]
for i, (title, desc, color) in enumerate(steps):
    x = 0.6 + i * 2.55
    add_rounded_shape(s, x, 1.4, 2.1, 0.55, color, title, 10, BG if color==BLUE else TITLE, True)
    add_text(s, x, 2.05, 2.1, 0.8, desc, 9, TEXT, align=PP_ALIGN.CENTER)
    if i < 4:
        add_arrow(s, x+2.2, 1.55, 0.25, 0.2)

add_text(s, 0.8, 3.2, 11.5, 0.3, '▎关键设计决策', 14, TITLE, True)
add_card(s, 0.8, 3.5, 3.8, 2.0, '为什么按帧分组而不是按行人？',
    ['SimpleLSTM 做法: 每条行人独立滑动窗口 → 丢失场景上下文',
     'Social-LSTM 做法: 同一帧的所有行人组成一个场景 → 保留社交交互',
     '场景窗口: 找到连续 20 帧内所有共存行人 → P 个行人一起预测',
     'P 从 1 (稀疏) 到 57 (拥挤) → 自然反映真实场景密度'],
    RED)

add_card(s, 5.0, 3.5, 3.8, 2.0, '归一化策略',
    ['标准做法: 以每个行人最后一帧观测位置为原点',
     'origin = obs[i, -1, :]  # (1, 2)',
     'obs_norm[i] = obs[i] - origin',
     'pred_norm[i] = pred[i] - origin',
     '优点: 模型学习相对位移，而非绝对坐标',
     '  不同场景的坐标范围差异很大 (数十米) → 归一化后统一尺度'],
    BLUE)

add_card(s, 9.2, 3.5, 3.5, 2.0, 'stride 参数',
    ['stride=1: 每帧都取窗口 (21,000+ 窗口) — 最全但慢',
     'stride=2: 隔帧采样 (10,600 窗口) — 推荐',
     'stride=4: 1/4 数据 (5,300 窗口) — 快速验证',
     'stride=8: 1/8 数据 (2,668 窗口) — 试水',
     '相邻窗口 95% 数据重叠 → stride=2 性价比最高',
     '当前选择: stride=1 (全量)'],
    ORANGE)

add_page_number(s, 5)

# ══════════════════════════════════════════════════════════════
# Slide 6 — Social-LSTM 架构 (1)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '04  模型架构：Social-LSTM (上)', 28, TITLE, True)
add_section_line(s, 1.0)

add_text(s, 0.8, 1.2, 11.5, 0.3, '▎整体架构：编码器-解码器 + 社交池化', 14, TITLE, True)

# 架构图 — 四个模块
modules = [
    ('输入', 'Obs\n(P, 8, 2)\n所有行人\n8帧位置', LGRAY),
    ('Position\nEmbedding', 'Linear(2→64)\n+ ReLU\n(P, 64)', LBLUE),
    ('Social\nPooling ★', '4×4 空间网格\n邻居h求和\n(P, 4,4,128)', LBLUE),
    ('LSTMCell\n(128→128)', '输入128维\n=pos(64)+social(64)\n逐时间步迭代', LBLUE),
]
for i, (title, desc, color) in enumerate(modules):
    x = 0.8 + i * 3.15
    add_rounded_shape(s, x, 1.6, 2.0, 0.6, color, title, 11, BG if color==BLUE else TITLE, True)
    add_text(s, x, 2.3, 2.0, 1.0, desc, 9, TEXT, align=PP_ALIGN.CENTER)
    if i < 3:
        add_arrow(s, x+2.1, 1.8, 0.3, 0.2)

# 社交池化详解
add_text(s, 0.8, 3.5, 11.5, 0.3, '▎核心创新：社交池化网格 (Social Pooling Grid)', 15, RED, True)

add_card(s, 0.8, 3.9, 6.0, 3.2, '网格构建过程',
    ['1. 以目标行人 i 当前位置为中心',
     '2. 建立 32×32 (米/像素) 的方形邻域窗口',
     '3. 将邻域划分为 4×4 = 16 个网格单元 (每单元 8×8)',
     '4. 对每个邻居 j：',
     '    a. 计算相对位置: (dx, dy) = (xⱼ−xᵢ, yⱼ−yᵢ)',
     '    b. 检查是否在邻域内: |dx|<16 且 |dy|<16',
     '    c. 映射到网格单元: cell_x = ⌊(dx+16)/32×4⌋, cell_y = ⌊(dy+16)/32×4⌋',
     '5. 对每个网格单元 (m, n)，求和落入该单元的所有邻居隐藏状态:',
     '    Hᵢ(m, n) = Σⱼ₌₁ᴺ 1_{mn}(xⱼ−xᵢ, yⱼ−yᵢ) · hⱼ',
     '6. 展平 → 社交嵌入: Linear(4×4×128=2048 → 64)',
     '',
     '★ 关键特性: 每时间步重新计算池化 (行人动态移动)  |  多邻居同单元→求和  |  空单元→零'],
    RED)

add_card(s, 7.2, 3.9, 5.6, 3.2, '为什么用网格池化？',
    ['▸ 空间局部性: 只有近邻 (32 单位内) 才相互影响',
     '  远处行人的行为通常无关',
     '',
     '▸ 排列不变性: 网格对邻居的排列不敏感',
     '  网格(2,3)有2个邻居 → 求和 → 同样的结果',
     '  无论输入顺序如何，池化结果一致',
     '',
     '▸ 计算高效: P 个行人 → O(P²) 对 → 16个matmul',
     '  P=75 时仅需 ~0.5ms (CPU)',
     '',
     '▸ 无需手动特征工程: 模型自动学习',
     '  哪些邻居重要、空间关系如何影响交互',
     '',
     '场景隔离: 打包训练时，不同场景行人之间互不交互',
     'scene_ids[i]==scene_ids[j] → 才能进入同一网格'],
    BLUE)

add_page_number(s, 6)

# ══════════════════════════════════════════════════════════════
# Slide 7 — Social-LSTM 架构 (2)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '04  模型架构：Social-LSTM (下)', 28, TITLE, True)
add_section_line(s, 1.0)

add_text(s, 0.8, 1.2, 11.5, 0.3, '▎逐时间步 LSTM 计算流程', 14, TITLE, True)

# 编码 vs 解码对比
add_card(s, 0.8, 1.6, 5.8, 2.5, '编码阶段 (Observation Encoding)',
    ['时间步 t = 0, 1, 2, ..., 7   (共 8 步)',
     '─── 每个时间步 ───',
     '① pos_t = obs[:, t, :]                         # 当前帧所有行人位置 (P, 2)',
     '② pos_emb = ReLU(Linear(pos_t))                # 位置嵌入 (P, 64)',
     '③ pooled = SocialPool(pos_t, h_{t-1})           # 社交池化 (P, 4, 4, 128)',
     '    (t=0 时 h 全零 → 池化结果全零 → 第1帧无社交信息)',
     '④ pooled_flat = pooled.view(P, -1)              # 展平 (P, 2048)',
     '⑤ social_emb = ReLU(Linear(pooled_flat))        # 社交嵌入 (P, 64)',
     '⑥ lstm_input = cat[pos_emb, social_emb]         # 拼接 (P, 128)',
     '⑦ h_t, c_t = LSTMCell(lstm_input, (h_{t-1}, c_{t-1}))  # 更新隐藏状态',
     '─── 编码结束 → h_7, c_7 包含 8 帧时空信息 ───'],
    BLUE)

add_card(s, 7.0, 1.6, 5.8, 2.5, '解码阶段 (Autoregressive Decoding)',
    ['时间步 t = 8, 9, ..., 19  (共 12 步)',
     '─── 每个时间步 ───',
     '① current_pos = current_pos + displacement       # 从预测位置出发',
     '   (第 8 步: current_pos = obs[:,-1,:] 即最后一帧观测位置)',
     '② pos_emb = ReLU(Linear(current_pos))            # (P, 64)',
     '③ pooled = SocialPool(current_pos, h_{t-1})       # 用预测位置做社交池化',
     '④ social_emb = ReLU(Linear(pooled_flat))          # (P, 64)',
     '⑤ lstm_input = cat[pos_emb, social_emb]            # (P, 128)',
     '⑥ h_t, c_t = LSTMCell(lstm_input, (h_{t-1}, c_{t-1}))',
     '⑦ displacement = Linear(h_t)                      # 预测位移 (P, 2)',
     '⑧ current_pos += displacement                     # 更新位置 → 下一步输入',
     '─── 堆叠 12 个 displacement → (P, 12, 2) 输出 ───'],
    GREEN)

# 关键参数
add_text(s, 0.8, 4.5, 11.5, 0.3, '▎为什么用 LSTMCell 而不是 nn.LSTM？', 14, TITLE, True)
add_card(s, 0.8, 4.8, 5.8, 1.5,
    'nn.LSTM (batch_first=True)',
    ['一次前向处理全部 8 帧 → 只得到最终隐藏状态',
     '无法在中间时间步插入社交池化',
     'SimpleLSTM 使用此方式 → 无社交交互 → 每条轨迹独立预测'],
    RED)
add_card(s, 7.0, 4.8, 5.8, 1.5,
    'nn.LSTMCell (逐时间步手动迭代)',
    ['每个时间步手动调用 → 可以在步间插入社交池化',
     '第 t 步的社交池化需要第 t-1 步的隐藏状态 h_{t-1}',
     '完全控制数据流 → 编码-解码可无缝衔接'],
    GREEN)

add_text(s, 0.8, 6.7, 11.5, 0.4, '参数量统计: PosEmbed(2→64) + SocialEmbed(2048→64) + LSTMCell(128→128) + Output(128→2) = 263,682 参数',
    size=11, color=META)
add_page_number(s, 7)

# ══════════════════════════════════════════════════════════════
# Slide 8 — 训练策略
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '05  训练策略与加速优化', 28, TITLE, True)
add_section_line(s, 1.0)

add_text(s, 0.8, 1.2, 11.5, 0.3, '▎核心瓶颈：为什么 CPU 利用率只有 50%？', 14, RED, True)

add_card(s, 0.8, 1.6, 5.8, 1.6, '问题诊断',
    ['逐场景训练: 每次只处理一个场景 → 平均只有 8.7 个行人',
     '矩阵运算: 8.7 × 128 → 太小, 无法充分利用 8 核 CPU',
     'LSTMCell 20 个时间步 → 串行执行, 无法并行',
     '结果: CPU 8个核心中只有 3-4 个在工作 → 利用率 ~50%',
     '20 步 LSTMCell 中每步都做 16 次池化 matmul → 160 次小矩阵乘法'],
    RED)

add_card(s, 7.0, 1.6, 5.8, 1.6, '解决方案：场景打包批处理',
    ['将 6 个场景拼成一批: 8.7人 × 6 = 52 人/批',
     '矩阵运算: 52 × 128 → 增大 6 倍, CPU 核心利用率大幅提升',
     '引入 scene_ids 参数: 确保不同场景行人之间互不交互',
     'same_scene = (scene_ids[i] == scene_ids[j]) → 掩码过滤',
     '验证: 打包推理结果与逐场景完全一致 (差异 < 1e-7)',
     '效果: 3.4x 加速 (200 场景: 3.5s → 1.0s)'],
    BLUE)

add_text(s, 0.8, 3.5, 11.5, 0.3, '▎完整训练配置', 14, TITLE, True)

# 训练配置表
config_items = [
    ['损失函数', 'MSE Loss', 'L = (1/N) Σ (ŷ − y)² , 在归一化坐标下计算'],
    ['优化器', 'Adam', 'lr=0.001, 默认 β=(0.9, 0.999)'],
    ['学习率调度', 'ReduceLROnPlateau', 'mode=min, factor=0.5, patience=10'],
    ['梯度管理', 'ClipGradNorm(max_norm=1.0)', '防止 LSTMCell 20步回传梯度爆炸'],
    ['累积步数', '1 (每批立即更新)', 'SCENE_BATCH=6 场景/批, 每批梯度已足够稳定'],
    ['Epochs', '100', '前30 epoch ADE快速下降, 50-100 epoch 精细调优'],
    ['验证策略', '每epoch验证', '取最佳 ADE 模型 (非最后 epoch 模型)'],
    ['训练设备', 'AMD 8845HS CPU', '或 Google Colab T4 GPU (免费)'],
]
for i, (k, v, desc) in enumerate(config_items):
    y = 3.9 + i * 0.38
    add_text(s, 0.8, y, 1.8, 0.35, k, 11, BLUE, True)
    add_text(s, 2.7, y, 2.5, 0.35, v, 11, TITLE, True)
    add_text(s, 5.3, y, 7.5, 0.35, desc, 10, META)

add_page_number(s, 8)

# ══════════════════════════════════════════════════════════════
# Slide 9 — 实验结果 (1)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '06  实验结果与分析 (上)', 28, TITLE, True)
add_section_line(s, 1.0)

add_text(s, 0.8, 1.2, 11.5, 0.3, '▎Social-LSTM 30 Epoch 试水训练 (stride=8, 1/8数据)', 14, TITLE, True)

# 训练曲线数据表
add_data_table(s, 0.8, 1.6,
    ['Epoch', '1', '5', '10', '15', '20', '25', '30'],
    [['Loss', '0.1135', '0.0497', '0.0403', '0.0335', '0.0278', '0.0249', '0.0214'],
     ['ADE(m)','0.747', '0.573', '0.535', '0.517', '0.534', '0.498', '0.483'],
     ['FDE(m)','1.433', '1.127', '1.053', '1.008', '1.019', '0.967', '0.924']],
    [0.9, 1.3, 1.3, 1.3, 1.3, 1.3, 1.3, 1.3],
    highlight_row=1)

add_text(s, 0.8, 3.0, 11.5, 0.3, '▎核心观察', 14, TITLE, True)
add_card(s, 0.8, 3.3, 3.8, 1.5, 'Loss 持续下降',
    ['0.114 → 0.021 (下降 5.4x)',
     '第 30 epoch 仍在下降',
     '未到达平台期 → 更多 epoch 有益',
     '梯度流畅, 无消失/爆炸'],
    GREEN)
add_card(s, 5.0, 3.3, 3.8, 1.5, 'ADE 稳步改善',
    ['0.747m → 0.482m (↓35%)',
     '初期收敛快 (前5轮 ↓0.17m)',
     '后期精细调优 (25→30 ↓0.01m)',
     '仅用 1/8 数据 + 30 epoch'],
    BLUE)
add_card(s, 9.2, 3.3, 3.5, 1.5, 'FDE 同步收敛',
    ['1.433m → 0.924m',
     'FDE 通常 > ADE (更难)',
     '最后一帧预测误差最大',
     '需要更多训练来降低远距误差'],
    ORANGE)

add_text(s, 0.8, 5.2, 11.5, 0.3, '▎模型对比 (ADE ↓ 越低越好)', 14, TITLE, True)
add_data_table(s, 0.8, 5.6,
    ['模型', '参数量', '训练Epoch', '数据量', 'ADE (m)', 'FDE (m)', '备注'],
    [['匀速基线 (CV)',  '-',     '-',  '-',    '0.480', '0.940', '常数速度外推: pred=obs[-1]+vel×t'],
     ['SimpleLSTM',     '5,400', '100','全量', '0.370', '0.720', '无社交交互, 单轨迹LSTM'],
     ['Social-LSTM',    '264K',  '30', '1/8',  '0.482', '0.924', '本次试水 ★ (未完整训练)'],
     ['Social-LSTM',    '264K',  '100','全量', '?',     '?',     '预期 < 0.37 (完整训练进行中)']],
    [1.0, 0.8, 0.9, 0.7, 1.0, 1.0, 2.0],
    highlight_row=3)
add_page_number(s, 9)

# ══════════════════════════════════════════════════════════════
# Slide 10 — 实验结果 (2)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '06  实验结果与分析 (下)', 28, TITLE, True)
add_section_line(s, 1.0)

add_text(s, 0.8, 1.2, 11.5, 0.3, '▎加速效果验证', 14, TITLE, True)
add_data_table(s, 0.8, 1.6,
    ['方案', '每批行人', '200场景耗时', '加速比', 'CPU利用率', '100ep预估'],
    [['逐场景 (原方案)',      '8.7',  '3.5s',  '1.0x', '~35%', '~18 小时'],
     ['场景打包×4',           '34.8', '1.7s',  '2.1x', '~55%', '~8 小时'],
     ['场景打包×6 (当前)',    '52.2', '1.0s',  '3.4x', '~70%', '~5 小时'],
     ['场景打包×8',           '69.6', '0.9s',  '3.9x', '~80%', '~4.5 小时']],
    [1.6, 1.1, 1.2, 1.0, 1.1, 1.2], highlight_row=2)

add_card(s, 0.8, 3.3, 5.8, 1.6, '场景打包正确性验证',
    ['验证方法: 同一份数据, 分两次推理 → 比较输出',
     '设置: 场景A(3人) + 场景B(2人) → 拼接(5人) + scene_ids',
     '逐场景: out_A = model(obs_A), out_B = model(obs_B)',
     '打包:   out_packed = model(obs_AB, scene_ids)',
     '结果: max|out_packed[:3] - out_A| = 2.98e-08',
     '      max|out_packed[3:] - out_B| = 8.94e-08',
     '结论: 打包推理与逐场景完全等价 (浮点精度范围) ✓'],
    GREEN)

add_card(s, 7.0, 3.3, 5.8, 1.6, '消融分析 (基于 30 epoch 数据)',
    ['▸ 社交池化 vs 无池化 (SimpleLSTM):',
     '  30 epoch: 0.609m vs 0.600m → 尚无明显优势',
     '  原因: Social-LSTM 参数量大12倍 → 收敛需要更多 epoch',
     '  预期: 完整100ep训练后反超 (论文结论 0.48→0.30)',
     '',
     '▸ 不同场景密度的影响:',
     '  Univ (75人/帧): 社交池化收益最大 (密集交互)',
     '  ETH (6人/帧):  社交信号弱, 与 SimpleLSTM 接近'],
    BLUE)

add_text(s, 0.8, 5.3, 11.5, 0.3, '▎与 SOTA 对比 (预期目标, 待完整训练后更新)', 14, TITLE, True)
add_data_table(s, 0.8, 5.7,
    ['方法', 'ETH', 'Hotel', 'Univ', 'Zara1', 'Zara2', '平均ADE'],
    [['Linear',        '1.33', '0.39', '0.82', '0.62', '0.52', '0.74'],
     ['Vanilla LSTM',  '1.09', '0.38', '0.61', '0.41', '0.52', '0.60'],
     ['Social-LSTM',   '1.09', '0.35', '0.50', '0.30', '0.37', '0.52'],
     ['Our SimpleLSTM','0.60', '—',    '—',    '—',    '—',    '— (单场景)'],
     ['Our SocialLSTM','?',    '?',    '?',    '?',    '?',    '? (待完整训练)']],
    [1.2, 1.0, 1.0, 1.0, 1.0, 1.0, 1.2], highlight_row=4)
add_text(s, 0.8, 7.0, 8, 0.3, '* SOTA 数据来自 Alahi et al. CVPR 2016 论文, 单位为米', 9, META)
add_page_number(s, 10)

# ══════════════════════════════════════════════════════════════
# Slide 11 — 代码架构
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '07  代码工程架构', 28, TITLE, True)
add_section_line(s, 1.0)

files = [
    ('src/config.py',            '全局超参数: 维度/学习率/epoch/路径 (单点修改, 全局生效)'),
    ('src/data_loader.py',       '原始数据加载: .txt → dict → 单轨迹窗口 (SimpleLSTM 基线用)'),
    ('src/scene_data_loader.py', '★ 场景级加载: 按帧分组 → 共存行人窗口 → 归一化 → PyTorch Dataset'),
    ('src/social_pooling.py',    '★ 社交池化模块: 4×4网格 + 场景隔离 + 邻居h求和 (核心创新代码)'),
    ('src/model.py',             '★ SocialLSTM + SimpleLSTM: LSTMCell逐步迭代 + 自回归解码'),
    ('src/train.py',             '训练入口: 场景打包批处理 + 梯度累积/裁剪 + 自动保存最佳模型'),
    ('src/evaluate.py',          '评估脚本: 同场景/跨场景(Leave-One-Out)/可视化 三模式'),
    ('src/quick_validate.py',    '快速验证: SimpleLSTM vs SocialLSTM 公平对比 (20 epoch)'),
    ('src/inference.py',         '单轨迹推理 + Matplotlib 可视化'),
    ('notebooks/train_on_colab.ipynb', 'Colab 训练: 7步一键运行, 自动下载模型'),
]
for i, (fname, desc) in enumerate(files):
    y = 1.3 + i * 0.52
    is_core = '★' in desc
    color = RED if is_core else BLUE
    add_text(s, 1.0, y, 3.2, 0.35, fname, 12, color, True)
    add_text(s, 4.3, y, 8.5, 0.35, desc, 11, TEXT)

add_text(s, 0.8, 7.0, 11.5, 0.3, '技术栈: Python 3.12 + PyTorch 2.11 + NumPy + Matplotlib + python-pptx  |  项目地址: github.com/confidentismylife/LonlyStydue',
    size=10, color=META, align=PP_ALIGN.CENTER)
add_page_number(s, 11)

# ══════════════════════════════════════════════════════════════
# Slide 12 — 未来工作
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '08  未来工作与论文方向', 28, TITLE, True)
add_section_line(s, 1.0)

add_text(s, 0.8, 1.2, 11.5, 0.4, '▎五大优化方向（按优先级排列）', 15, TITLE, True)

add_card(s, 0.8, 1.7, 3.8, 2.5, '① 火灾场景特征融入 ★★',
    ['核心创新点 → 论文差异化关键',
     '引入 FDS 火灾动力学模拟数据:',
     '  烟雾密度场 / CO 浓度 / 温度场',
     '  出口方向向量场 (引导方向)',
     '  人群密度热力图 (局部拥堵)',
     '多模态输入:',
     '  lstm_input = [pos_emb, social_emb, fire_emb]',
     '  模型同时考虑人-人交互 & 人-环境交互',
     '贡献: 首个融合火灾环境特征的行人预测模型'],
    RED)

add_card(s, 5.0, 1.7, 3.8, 2.5, '② 损失函数升级',
    ['MSE → 双变量高斯负对数似然 (NLL)',
     '  模型输出: (μx, μy, σx, σy, ρ) 5 参数',
     '  L = -log P(y | μ, Σ)',
     '  学习预测的不确定性 σ → 更符合现实',
     '',
     'Variety Loss (多模态预测):',
     '  生成 K=20 条候选轨迹',
     '  取与真值最接近的轨迹算 loss',
     '  处理行人行为的固有不确定性',
     '  (向左绕 vs 向右绕 都合理)'],
    BLUE)

add_card(s, 9.2, 1.7, 3.5, 2.5, '③ 模型架构升级',
    ['注意力池化: 替代固定 4×4 网格',
     '  学习哪些邻居重要 (动态权重)',
     '  Attention(Q=hᵢ, K=hⱼ, V=hⱼ)',
     '',
     '图神经网络: GAT/GCN',
     '  行人 = 图节点, 交互 = 边',
     '  边的权重由相对距离决定',
     '',
     'Transformer / Perceiver',
     '  并行处理所有时间步',
     '  全局感受野'],
    GREEN)

add_card(s, 0.8, 4.5, 3.8, 1.8, '④ 工程部署',
    ['ONNX 导出模型 → 跨平台推理',
     'TensorRT / OpenVINO 量化加速',
     'C++ 运行时 (实时性 < 1ms)',
     '集成到数字孪生可视化平台',
     'Web 端实时预测演示'],
    ORANGE)

add_card(s, 5.0, 4.5, 3.8, 1.8, '⑤ 完整实验体系',
    ['Leave-One-Out 5折交叉验证',
     '消融: 社交池化/场景密度/帧数',
     '鲁棒性: 加噪声/遮挡/缺失帧',
     '火灾模拟场景定量评估',
     '与最新 SOTA (2024-25) 对比'],
    BLUE)

add_card(s, 9.2, 4.5, 3.5, 1.8, '论文投稿方向',
    ['Safety Science (安全科学)',
     'Fire Technology (消防技术)',
     'Physica A (统计力学)',
     'Building and Environment',
     'Automation in Construction'],
    TITLE)

add_page_number(s, 12)

# ══════════════════════════════════════════════════════════════
# Slide 13 — 总结
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.3, 10, 0.6, '总  结', 32, TITLE, True)
add_section_line(s, 1.0)

add_card(s, 0.8, 1.3, 3.8, 2.0, '✅ 已完成',
    ['完整数据管道 (ETH/UCY 5场景)',
     'Social-LSTM 模型完整实现',
     '  社交池化层 (4×4 网格)',
     '  LSTMCell 逐时间步推理',
     '  编码器-解码器自回归预测',
     '场景打包批处理 (3.4x 加速)',
     'SimpleLSTM 基线对比',
     '30 epoch 试水验证 (ADE=0.482m)',
     'Colab 云端训练支持'],
    GREEN)

add_card(s, 5.0, 1.3, 3.8, 2.0, '🔄 进行中',
    ['100 epoch 完整训练 (stride=1)',
     'Leave-One-Out 交叉验证',
     '超参数网格搜索',
     '与 SOTA 定量对比',
     '训练曲线 + 可视化完善'],
    ORANGE)

add_card(s, 9.2, 1.3, 3.5, 2.0, '📋 计划中',
    ['火灾环境特征集成 (FDS)',
     '双变量高斯 NLL 损失',
     '多模态轨迹预测 (K=20)',
     '注意力/Transformer 升级',
     'ONNX 部署 + 实时推理',
     '论 文 撰 写'],
    RED)

add_text(s, 0.8, 3.7, 11.5, 0.4, '▎核心贡献总结', 16, TITLE, True)

contribs = [
    '① 从零实现完整 Social-LSTM 模型: 社交池化网格 + LSTMCell 逐时间步 + 自回归解码，263K 参数可在消费级 CPU 上训练',
    '② 提出"场景打包批处理"策略: 多场景拼接 + scene_ids 隔离，CPU 训练速度提升 3.4 倍，从 ~18 小时降至 ~5 小时',
    '③ 建立完整评估体系: 支持 3 种评估模式 (同场景测试 / Leave-One-Out 跨场景 / 可视化对比)',
    '④ 为火灾疏散智能决策奠定基础: 行人轨迹预测 → 实时拥堵预警 → 智能疏散引导 → 减少伤亡',
]
for i, c in enumerate(contribs):
    add_text(s, 1.2, 4.2 + i * 0.45, 11, 0.4, c, 13, TEXT)

add_text(s, 0.8, 6.2, 11.5, 0.6, '谢  谢', 40, RED, True, PP_ALIGN.CENTER)
add_text(s, 0.8, 6.9, 11.5, 0.4, '欢迎提问与讨论', 16, META, align=PP_ALIGN.CENTER)
add_page_number(s, 13)

# ── Save ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           '组会演示_火灾疏散轨迹预测_v2.pptx')
prs.save(output_path)
print(f'OK → {output_path}')
